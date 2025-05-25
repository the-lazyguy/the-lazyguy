
"""
PowerPoint to PDF Summary Tool
"""

import os
import sys
import argparse
from pathlib import Path
from io import BytesIO
from typing import List, Tuple, Optional
import tempfile
import threading
import webbrowser
from datetime import datetime

# PowerPoint processing
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.lib.colors import black, darkblue, darkgray
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.platypus.tableofcontents import TableOfContents
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

# Image processing
from PIL import Image as PILImage

# Web interface
from flask import Flask, render_template_string, request, send_file, flash, redirect, url_for, jsonify
from werkzeug.utils import secure_filename


class PowerPointToPDFConverter:
    """Converts PowerPoint presentations to clean PDF summaries."""
    
    def __init__(self, output_path: str):
        self.output_path = output_path
        self.doc = SimpleDocTemplate(
            output_path,
            pagesize=A4,
            rightMargin=2*cm,
            leftMargin=2*cm,
            topMargin=2.5*cm,
            bottomMargin=2*cm
        )
        self.story = []
        self.styles = self._create_styles()
        self.image_counter = 0
        self.include_images = True  # Default to include images
        
    def _create_styles(self) -> dict:
        """Create custom styles for the PDF document."""
        styles = getSampleStyleSheet()
        
        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            spaceBefore=20,
            textColor=darkblue,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Slide heading style
        slide_heading_style = ParagraphStyle(
            'SlideHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            spaceBefore=20,
            textColor=darkblue,
            fontName='Helvetica-Bold',
            borderWidth=1,
            borderColor=darkgray,
            borderPadding=8,
            backColor='#f0f0f0'
        )
        
        # Content style
        content_style = ParagraphStyle(
            'Content',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=8,
            spaceBefore=4,
            textColor=black,
            alignment=TA_JUSTIFY,
            fontName='Helvetica',
            lineHeight=1.4
        )
        
        # Bullet point style
        bullet_style = ParagraphStyle(
            'BulletPoint',
            parent=content_style,
            leftIndent=20,
            bulletIndent=8,
            spaceAfter=4
        )
        
        return {
            'title': title_style,
            'slide_heading': slide_heading_style,
            'content': content_style,
            'bullet': bullet_style
        }
    
    def _extract_text_from_shape(self, shape) -> str:
        """Extract text content from a PowerPoint shape."""
        text_content = []
        
        if hasattr(shape, 'text') and shape.text.strip():
            text_content.append(shape.text.strip())
        
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
        
        return '\n'.join(text_content)
    
    def _extract_image_from_shape(self, shape, slide_num: int) -> Optional[str]:
        """Extract image from a PowerPoint shape and save it temporarily."""
        if not isinstance(shape, Picture):
            return None
        
        try:
            # Get image data
            image_stream = BytesIO(shape.image.blob)
            pil_image = PILImage.open(image_stream)
            
            # Resize image if too large
            max_width, max_height = 400, 300
            if pil_image.width > max_width or pil_image.height > max_height:
                pil_image.thumbnail((max_width, max_height), PILImage.Resampling.LANCZOS)
            
            # Save temporary image
            self.image_counter += 1
            temp_path = f"temp_image_{slide_num}_{self.image_counter}.png"
            pil_image.save(temp_path, "PNG")
            
            return temp_path
            
        except Exception as e:
            print(f"Warning: Could not extract image from slide {slide_num}: {e}")
            return None
    
    def _process_slide(self, slide, slide_num: int) -> Tuple[str, List[str], List[str]]:
        """Process a single slide and extract content."""
        slide_title = ""
        text_content = []
        images = []
        
        for shape in slide.shapes:
            # Handle grouped shapes
            if isinstance(shape, GroupShape):
                for grouped_shape in shape.shapes:
                    text = self._extract_text_from_shape(grouped_shape)
                    if text:
                        text_content.append(text)
            else:
                # Extract text
                text = self._extract_text_from_shape(shape)
                if text:
                    # First significant text becomes the slide title
                    if not slide_title and len(text) < 100:
                        slide_title = text
                    else:
                        text_content.append(text)
                
                # Extract images
                image_path = self._extract_image_from_shape(shape, slide_num)
                if image_path:
                    images.append(image_path)
        
        return slide_title, text_content, images
    
    def _clean_text(self, text: str) -> str:
        """Clean and format text content."""
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Handle common PowerPoint artifacts
        text = text.replace('', '').replace('', '')
        
        return text
    
    def _add_slide_content(self, slide_num: int, title: str, content: List[str], images: List[str]):
        """Add slide content to the PDF story."""
        # Slide heading
        if title:
            heading_text = f"Slide {slide_num}: {self._clean_text(title)}"
        else:
            heading_text = f"Slide {slide_num}"
        
        self.story.append(Paragraph(heading_text, self.styles['slide_heading']))
        self.story.append(Spacer(1, 0.2*inch))
        
        # Content
        for text in content:
            cleaned_text = self._clean_text(text)
            if not cleaned_text:
                continue
                
            # Check if text looks like bullet points
            if any(cleaned_text.startswith(marker) for marker in ['•', '-', '*', '○']):
                # Split into bullet points
                lines = cleaned_text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        # Remove bullet markers and add consistent formatting
                        line = line.lstrip('•-*○ ').strip()
                        if line:
                            bullet_text = f"• {line}"
                            self.story.append(Paragraph(bullet_text, self.styles['bullet']))
            else:
                # Regular paragraph
                self.story.append(Paragraph(cleaned_text, self.styles['content']))
        
        # Add images
        if self.include_images:
            for image_path in images:
                try:
                    # Add some space before image
                    self.story.append(Spacer(1, 0.1*inch))
                    
                    # Create image with appropriate sizing
                    img = Image(image_path, width=4*inch, height=3*inch, kind='proportional')
                    self.story.append(img)
                    
                    # Add space after image
                    self.story.append(Spacer(1, 0.1*inch))
                    
                except Exception as e:
                    print(f"Warning: Could not add image {image_path}: {e}")
        
        # Add space between slides
        self.story.append(Spacer(1, 0.3*inch))
    
    def convert(self, ppt_path: str, title: Optional[str] = None):
        """Convert PowerPoint to PDF summary."""
        try:
            # Load presentation
            prs = Presentation(ppt_path)
            
            # Add title page
            if not title:
                title = f"Summary of {Path(ppt_path).stem}"
            
            self.story.append(Paragraph(title, self.styles['title']))
            self.story.append(Spacer(1, 0.5*inch))
            
            # Add summary info
            summary_info = f"Generated from: {Path(ppt_path).name}<br/>Total slides: {len(prs.slides)}<br/>Created: {Path(ppt_path).stat().st_mtime}"
            self.story.append(Paragraph(summary_info, self.styles['content']))
            self.story.append(PageBreak())
            
            # Process each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                print(f"Processing slide {slide_num}...")
                
                slide_title, content, images = self._process_slide(slide, slide_num)
                self._add_slide_content(slide_num, slide_title, content, images)
            
            # Build PDF
            print("Generating PDF...")
            self.doc.build(self.story)
            print(f"PDF generated successfully: {self.output_path}")
            
        except Exception as e:
            print(f"Error converting PowerPoint: {e}")
            raise
        
        finally:
            # Clean up temporary image files
            self._cleanup_temp_files()
    
    def _cleanup_temp_files(self):
        """Remove temporary image files."""
        for i in range(1, self.image_counter + 1):
            for slide_num in range(1, 100):  # Reasonable upper bound
                temp_file = f"temp_image_{slide_num}_{i}.png"
                if os.path.exists(temp_file):
                    try:
                        os.remove(temp_file)
                    except:
                        pass


def main():
    """Main function to handle command line usage."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentations to clean PDF summaries"
    )
    parser.add_argument("--web", action="store_true", help="Launch web interface")
    parser.add_argument("--cli", action="store_true", help="Force command line mode")
    parser.add_argument("input", nargs="?", help="Input PowerPoint file (.pptx)")
    parser.add_argument("output", nargs="?", help="Output PDF file")
    parser.add_argument("--title", help="Custom title for the PDF", default=None)
    parser.add_argument("--port", type=int, default=5000, help="Port for web interface")
    
    args = parser.parse_args()
    
    # If --web is specified or no arguments provided, launch web interface
    if args.web or (not args.cli and not args.input and not args.output):
        launch_web_interface(args.port)
        return
    
    # Command line mode
    if not args.input or not args.output:
        print("Error: Input and output files are required for command line mode.")
        print("Usage:")
        print("  python ppt_to_pdf.py input.pptx output.pdf [--title 'Custom Title']")
        print("  python ppt_to_pdf.py --web [--port 5000]")
        print("  python ppt_to_pdf.py (launches web interface by default)")
        sys.exit(1)
    
    # Validate input file
    if not os.path.exists(args.input):
        print(f"Error: Input file '{args.input}' not found.")
        sys.exit(1)
    
    if not args.input.lower().endswith(('.pptx', '.ppt')):
        print("Warning: Input file should be a PowerPoint presentation (.pptx or .ppt)")
    
    # Create output directory if needed
    output_dir = os.path.dirname(args.output)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        # Convert presentation
        converter = PowerPointToPDFConverter(args.output)
        converter.convert(args.input, args.title)
        
    except Exception as e:
        print(f"Conversion failed: {e}")
        sys.exit(1)



# Web Interface
class WebInterface:
    """Flask web interface for the PowerPoint to PDF converter."""
    
    def __init__(self):
        self.app = Flask(__name__)
        self.app.secret_key = 'ppt_to_pdf_converter_2024'
        self.app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size
        self.upload_folder = tempfile.mkdtemp()
        self.output_folder = tempfile.mkdtemp()
        
        # Setup routes
        self.setup_routes()
    
    def setup_routes(self):
        """Setup Flask routes."""
        
        @self.app.route('/')
        def index():
            return render_template_string(self.get_html_template())
        
        @self.app.route('/upload', methods=['POST'])
        def upload_file():
            try:
                if 'file' not in request.files:
                    return jsonify({'error': 'No file selected'}), 400
                
                file = request.files['file']
                if file.filename == '':
                    return jsonify({'error': 'No file selected'}), 400
                
                if not file.filename.lower().endswith(('.pptx', '.ppt')):
                    return jsonify({'error': 'Please upload a PowerPoint file (.pptx or .ppt)'}), 400
                
                # Save uploaded file
                filename = secure_filename(file.filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                unique_filename = f"{timestamp}_{filename}"
                file_path = os.path.join(self.upload_folder, unique_filename)
                file.save(file_path)
                
                # Get conversion options
                custom_title = request.form.get('title', '').strip()
                include_images = request.form.get('include_images') == 'on'
                
                # Generate output filename
                output_filename = f"{Path(filename).stem}_summary.pdf"
                output_path = os.path.join(self.output_folder, f"{timestamp}_{output_filename}")
                
                # Convert to PDF
                converter = PowerPointToPDFConverter(output_path)
                converter.include_images = include_images
                converter.convert(file_path, custom_title if custom_title else None)
                
                # Clean up input file
                os.remove(file_path)
                
                return jsonify({
                    'success': True, 
                    'download_url': f'/download/{timestamp}_{output_filename}',
                    'filename': output_filename
                })
                
            except Exception as e:
                return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
        
        @self.app.route('/download/<filename>')
        def download_file(filename):
            try:
                file_path = os.path.join(self.output_folder, filename)
                if os.path.exists(file_path):
                    return send_file(file_path, as_attachment=True, 
                                   download_name=filename.split('_', 1)[1] if '_' in filename else filename)
                else:
                    return "File not found", 404
            except Exception as e:
                return f"Error downloading file: {str(e)}", 500
    
    def get_html_template(self):
        """Return the HTML template for the web interface."""
        return '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint to PDF Converter</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 20px;
        }
        
        .container {
            background: white;
            border-radius: 20px;
            box-shadow: 0 20px 40px rgba(0,0,0,0.1);
            padding: 40px;
            max-width: 600px;
            width: 100%;
        }
        
        .header {
            text-align: center;
            margin-bottom: 40px;
        }
        
        .header h1 {
            color: #333;
            font-size: 2.5em;
            margin-bottom: 10px;
            font-weight: 300;
        }
        
        .header p {
            color: #666;
            font-size: 1.1em;
            line-height: 1.6;
        }
        
        .upload-section {
            margin-bottom: 30px;
        }
        
        .upload-area {
            border: 3px dashed #ddd;
            border-radius: 15px;
            padding: 40px 20px;
            text-align: center;
            background: #fafafa;
            transition: all 0.3s ease;
            cursor: pointer;
        }
        
        .upload-area:hover {
            border-color: #667eea;
            background: #f0f4ff;
        }
        
        .upload-area.dragover {
            border-color: #667eea;
            background: #e8f0ff;
            transform: scale(1.02);
        }
        
        .upload-icon {
            font-size: 3em;
            color: #667eea;
            margin-bottom: 20px;
        }
        
        .upload-text {
            color: #666;
            font-size: 1.1em;
            margin-bottom: 15px;
        }
        
        .file-input {
            display: none;
        }
        
        .btn {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            border: none;
            padding: 12px 30px;
            border-radius: 25px;
            font-size: 1em;
            cursor: pointer;
            transition: all 0.3s ease;
            text-decoration: none;
            display: inline-block;
        }
        
        .btn:hover {
            transform: translateY(-2px);
            box-shadow: 0 10px 20px rgba(0,0,0,0.2);
        }
        
        .btn:disabled {
            opacity: 0.6;
            cursor: not-allowed;
            transform: none;
        }
        
        .options {
            background: #f8f9fa;
            padding: 25px;
            border-radius: 15px;
            margin-bottom: 30px;
        }
        
        .option-group {
            margin-bottom: 20px;
        }
        
        .option-group:last-child {
            margin-bottom: 0;
        }
        
        .option-group label {
            display: block;
            margin-bottom: 8px;
            font-weight: 500;
            color: #333;
        }
        
        .option-group input[type="text"] {
            width: 100%;
            padding: 12px 15px;
            border: 2px solid #e9ecef;
            border-radius: 10px;
            font-size: 1em;
            transition: border-color 0.3s ease;
        }
        
        .option-group input[type="text"]:focus {
            outline: none;
            border-color: #667eea;
        }
        
        .checkbox-group {
            display: flex;
            align-items: center;
            gap: 10px;
        }
        
        .checkbox-group input[type="checkbox"] {
            width: 18px;
            height: 18px;
            accent-color: #667eea;
        }
        
        .progress {
            display: none;
            margin: 20px 0;
        }
        
        .progress-bar {
            width: 100%;
            height: 8px;
            background: #e9ecef;
            border-radius: 10px;
            overflow: hidden;
        }
        
        .progress-fill {
            height: 100%;
            background: linear-gradient(90deg, #667eea, #764ba2);
            border-radius: 10px;
            animation: progress-animation 2s ease-in-out infinite;
        }
        
        @keyframes progress-animation {
            0% { width: 0%; }
            50% { width: 100%; }
            100% { width: 0%; }
        }
        
        .result {
            display: none;
            text-align: center;
            padding: 30px;
            background: #d4edda;
            border: 1px solid #c3e6cb;
            border-radius: 15px;
            margin-top: 20px;
        }
        
        .result.error {
            background: #f8d7da;
            border-color: #f5c6cb;
        }
        
        .result-icon {
            font-size: 3em;
            margin-bottom: 15px;
        }
        
        .success .result-icon {
            color: #28a745;
        }
        
        .error .result-icon {
            color: #dc3545;
        }
        
        .selected-file {
            display: none;
            background: #e8f0ff;
            padding: 15px;
            border-radius: 10px;
            margin-top: 15px;
            border: 1px solid #667eea;
        }
        
        .file-info {
            display: flex;
            align-items: center;
            gap: 10px;
            color: #333;
        }
        
        .file-icon {
            color: #667eea;
            font-size: 1.5em;
        }
        
        .convert-section {
            text-align: center;
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>📊 PowerPoint to PDF</h1>
            <p>Convert your PowerPoint presentations into clean, organized PDF summaries with just a few clicks.</p>
        </div>
        
        <form id="uploadForm" enctype="multipart/form-data">
            <div class="upload-section">
                <div class="upload-area" id="uploadArea">
                    <div class="upload-icon">📁</div>
                    <div class="upload-text">
                        <strong>Click to select</strong> or drag and drop your PowerPoint file here
                    </div>
                    <div style="color: #999; font-size: 0.9em; margin-top: 10px;">
                        Supported formats: .pptx, .ppt (Max size: 100MB)
                    </div>
                </div>
                <input type="file" id="fileInput" name="file" class="file-input" accept=".pptx,.ppt">
                
                <div class="selected-file" id="selectedFile">
                    <div class="file-info">
                        <span class="file-icon">📄</span>
                        <span id="fileName"></span>
                    </div>
                </div>
            </div>
            
            <div class="options">
                <div class="option-group">
                    <label for="titleInput">Custom Title (optional):</label>
                    <input type="text" id="titleInput" name="title" placeholder="Enter a custom title for your PDF summary">
                </div>
                
                <div class="option-group">
                    <div class="checkbox-group">
                        <input type="checkbox" id="includeImages" name="include_images" checked>
                        <label for="includeImages">Include images from slides</label>
                    </div>
                </div>
            </div>
            
            <div class="convert-section">
                <button type="submit" class="btn" id="convertBtn">
                    🔄 Convert to PDF
                </button>
            </div>
        </form>
        
        <div class="progress" id="progress">
            <div style="text-align: center; margin-bottom: 15px; color: #666;">
                Converting your presentation...
            </div>
            <div class="progress-bar">
                <div class="progress-fill"></div>
            </div>
        </div>
        
        <div class="result" id="result">
            <div class="result-icon">✅</div>
            <div id="resultMessage"></div>
            <div style="margin-top: 20px;">
                <a href="#" id="downloadBtn" class="btn">📥 Download PDF</a>
            </div>
        </div>
    </div>
    
    <script>
        const uploadArea = document.getElementById('uploadArea');
        const fileInput = document.getElementById('fileInput');
        const uploadForm = document.getElementById('uploadForm');
        const selectedFile = document.getElementById('selectedFile');
        const fileName = document.getElementById('fileName');
        const convertBtn = document.getElementById('convertBtn');
        const progress = document.getElementById('progress');
        const result = document.getElementById('result');
        const resultMessage = document.getElementById('resultMessage');
        const downloadBtn = document.getElementById('downloadBtn');
        
        // File upload handling
        uploadArea.addEventListener('click', () => fileInput.click());
        
        uploadArea.addEventListener('dragover', (e) => {
            e.preventDefault();
            uploadArea.classList.add('dragover');
        });
        
        uploadArea.addEventListener('dragleave', () => {
            uploadArea.classList.remove('dragover');
        });
        
        uploadArea.addEventListener('drop', (e) => {
            e.preventDefault();
            uploadArea.classList.remove('dragover');
            
            const files = e.dataTransfer.files;
            if (files.length > 0) {
                fileInput.files = files;
                showSelectedFile(files[0]);
            }
        });
        
        fileInput.addEventListener('change', (e) => {
            if (e.target.files.length > 0) {
                showSelectedFile(e.target.files[0]);
            }
        });
        
        function showSelectedFile(file) {
            fileName.textContent = file.name;
            selectedFile.style.display = 'block';
            convertBtn.disabled = false;
        }
        
        // Form submission
        uploadForm.addEventListener('submit', async (e) => {
            e.preventDefault();
            
            if (!fileInput.files.length) {
                alert('Please select a PowerPoint file first.');
                return;
            }
            
            const formData = new FormData(uploadForm);
            
            // Show progress
            convertBtn.disabled = true;
            progress.style.display = 'block';
            result.style.display = 'none';
            
            try {
                const response = await fetch('/upload', {
                    method: 'POST',
                    body: formData
                });
                
                const data = await response.json();
                
                progress.style.display = 'none';
                result.style.display = 'block';
                
                if (data.success) {
                    result.className = 'result success';
                    result.querySelector('.result-icon').textContent = '✅';
                    resultMessage.textContent = `Your PDF summary has been generated successfully!`;
                    downloadBtn.href = data.download_url;
                    downloadBtn.style.display = 'inline-block';
                } else {
                    throw new Error(data.error || 'Conversion failed');
                }
                
            } catch (error) {
                progress.style.display = 'none';
                result.style.display = 'block';
                result.className = 'result error';
                result.querySelector('.result-icon').textContent = '❌';
                resultMessage.textContent = `Error: ${error.message}`;
                downloadBtn.style.display = 'none';
            }
            
            convertBtn.disabled = false;
        });
    </script>
</body>
</html>
        '''
    
    def run(self, port=5000, debug=False):
        """Run the Flask web application."""
        self.app.run(host='0.0.0.0', port=port, debug=debug)


def launch_web_interface(port=5000):
    """Launch the web interface."""
    web_app = WebInterface()
    
    print(f"\n🚀 PowerPoint to PDF Converter Web Interface")
    print(f"   Server starting on http://localhost:{port}")
    print(f"   Press Ctrl+C to stop the server\n")
    
    # Open browser automatically
    def open_browser():
        webbrowser.open(f'http://localhost:{port}')
    
    timer = threading.Timer(1.0, open_browser)
    timer.start()
    
    try:
        web_app.run(port=port)
    except KeyboardInterrupt:
        print("\n🛑 Server stopped by user")
    except Exception as e:
        print(f"\n❌ Error starting server: {e}")


if __name__ == "__main__":
    main()


